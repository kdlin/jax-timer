from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color tokens ────────────────────────────────────────────────
BG          = RGBColor(0x13, 0x13, 0x13)
SURFACE_LOW = RGBColor(0x1b, 0x1b, 0x1b)
SURFACE     = RGBColor(0x22, 0x22, 0x22)
SURFACE_HI  = RGBColor(0x2a, 0x2a, 0x2a)
SURFACE_HHI = RGBColor(0x35, 0x35, 0x35)
OUTLINE     = RGBColor(0x47, 0x47, 0x47)
PRIMARY     = RGBColor(0xFF, 0xFF, 0xFF)
ON_SURFACE  = RGBColor(0xC6, 0xC6, 0xC6)
TERTIARY    = RGBColor(0x6B, 0xFD, 0xAF)
BREAK_ACC   = RGBColor(0xF1, 0x9A, 0x8E)
ERROR       = RGBColor(0xFF, 0xB4, 0xAB)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]

def bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG

def box(slide, l, t, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Inter'
    return tb

def badge(slide, text, l, t, color=TERTIARY):
    w, h = 1.0, 0.28
    box(slide, l, t, w, h, SURFACE_HHI, OUTLINE)
    tb = slide.shapes.add_textbox(Inches(l+0.08), Inches(t+0.04), Inches(w-0.16), Inches(h-0.08))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(7)
    run.font.color.rgb = color
    run.font.bold = True
    run.font.name = 'Inter'

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)

# Faint ambient panels
box(s1, 9.5, -0.5, 7, 5, RGBColor(0x0D, 0x1F, 0x17))
box(s1, 0.0,  5.5, 6, 5, RGBColor(0x14, 0x14, 0x1E))

badge(s1, 'OBSIDIAN HUD', 0.6, 0.5)
txt(s1, 'Design System', 0.6, 0.85, 9, 1.5, 60, PRIMARY, bold=True)
txt(s1, 'v1.0', 0.65, 2.15, 2, 0.4, 13, ON_SURFACE)
txt(s1, 'High-End Minimalist Dark Mode\nInspired by surgical precision and editorial clarity.', 0.6, 2.6, 8, 0.9, 13, ON_SURFACE)

# HUD preview card
box(s1, 10.6, 1.2, 4.8, 6.4, SURFACE_LOW, OUTLINE)
txt(s1, 'FOCUS', 11.0, 1.55, 2, 0.3, 7, TERTIARY, bold=True)
txt(s1, '25:00', 10.8, 1.95, 4.8, 1.8, 64, PRIMARY, bold=True)
box(s1, 11.0, 4.2, 4.1, 0.05, SURFACE_HHI)
box(s1, 11.0, 4.2, 2.0, 0.05, PRIMARY)
txt(s1, '12:30', 11.0, 4.32, 1.5, 0.3, 8, ON_SURFACE)
txt(s1, '25:00', 14.3, 4.32, 1.0, 0.3, 8, ON_SURFACE, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Color Palette
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)

badge(s2, 'FOUNDATIONS', 0.6, 0.4)
txt(s2, 'Color Palette', 0.6, 0.78, 10, 0.8, 36, PRIMARY, bold=True)
txt(s2, 'Tonal layering — surfaces as physical layers of polished obsidian', 0.6, 1.52, 12, 0.4, 12, ON_SURFACE)

swatches = [
    ('#131313', 'surface',               BG,          PRIMARY),
    ('#1b1b1b', 'surface-container-low', SURFACE_LOW, PRIMARY),
    ('#2a2a2a', 'surface-container-high',SURFACE_HI,  PRIMARY),
    ('#353535', 'surface-container-highest', SURFACE_HHI, PRIMARY),
    ('#474747', 'outline-variant',       OUTLINE,     PRIMARY),
    ('#ffffff', 'primary',               PRIMARY,     BG),
    ('#6bfdaf', 'tertiary',              TERTIARY,    BG),
    ('#f19a8e', 'break-accent',          BREAK_ACC,   BG),
]

cols = 4
for i, (hex_val, name, color, text_color) in enumerate(swatches):
    col = i % cols
    row = i // cols
    lx = 0.5 + col * 3.8
    ty = 2.1 + row * 2.9
    box(s2, lx, ty, 3.5, 2.2, color, OUTLINE)
    txt(s2, hex_val.upper(), lx+0.15, ty+1.65, 3.0, 0.4, 9,  text_color, bold=True)
    txt(s2, name,             lx+0.15, ty+1.9,  3.0, 0.35, 8, text_color)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Typography
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)

badge(s3, 'FOUNDATIONS', 0.6, 0.4)
txt(s3, 'Typography', 0.6, 0.78, 10, 0.8, 36, PRIMARY, bold=True)
txt(s3, 'Inter — editorial hierarchy with drastic scale shifts', 0.6, 1.52, 12, 0.4, 12, ON_SURFACE)

type_rows = [
    ('Display',  '56pt  /  -0.02em tracking',  'The Obsidian HUD',         40, PRIMARY,   True,  0.95),
    ('Headline', '32pt  /  high contrast',      'Section Headline',         26, PRIMARY,   True,  0.78),
    ('Title',    '22pt  /  component headers',  'Component Title',          18, PRIMARY,   False, 0.60),
    ('Body',     '13pt  /  on-surface-variant', 'Body text for readability.',13, ON_SURFACE,False, 0.50),
    ('Label',    '9pt   /  uppercase + 0.05em', 'LABEL  •  BADGE  •  META', 9,  TERTIARY,  True,  0.40),
]

ty = 2.1
for role, size_label, sample, pt, color, bold, row_h in type_rows:
    box(s3, 0.5, ty, 15, 0.03, SURFACE_HI)
    txt(s3, role,       0.65, ty+0.08, 2.4, 0.45, 9,  ON_SURFACE)
    txt(s3, size_label, 3.2,  ty+0.08, 3.2, 0.45, 8,  OUTLINE)
    txt(s3, sample,     6.5,  ty+0.04, 9.0, row_h+0.1, pt, color, bold=bold)
    ty += row_h + 0.35

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Components
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)

badge(s4, 'COMPONENTS', 0.6, 0.4)
txt(s4, 'UI Components', 0.6, 0.78, 10, 0.8, 36, PRIMARY, bold=True)

# Buttons
txt(s4, 'BUTTONS', 0.6, 1.72, 4, 0.35, 8, ON_SURFACE, bold=True)
box(s4, 0.6, 2.12, 2.4, 0.58, PRIMARY, None)
txt(s4, 'Primary Action', 0.7, 2.22, 2.3, 0.38, 10, BG, bold=True)
box(s4, 3.2, 2.12, 2.2, 0.58, BG, OUTLINE)
txt(s4, 'Secondary', 3.32, 2.22, 2.1, 0.38, 10, PRIMARY)
box(s4, 5.6, 2.12, 1.8, 0.58, BG, None)
txt(s4, 'Tertiary', 5.7, 2.22, 1.7, 0.38, 10, ON_SURFACE)

# Badges
txt(s4, 'BADGES', 0.6, 3.05, 4, 0.35, 8, ON_SURFACE, bold=True)
for i, (label_txt, color) in enumerate([('FOCUS', TERTIARY), ('BREAK', BREAK_ACC), ('ACTIVE', TERTIARY), ('ERROR', ERROR), ('POST', TERTIARY)]):
    badge(s4, label_txt, 0.6 + i * 1.15, 3.45, color)

# HUD Card
txt(s4, 'FLOATING HUD  (glassmorphism)', 0.6, 4.15, 8, 0.35, 8, ON_SURFACE, bold=True)
box(s4, 0.6, 4.55, 5.6, 4.0, SURFACE_LOW, OUTLINE)
txt(s4, 'FOCUS', 0.88, 4.78, 2, 0.3, 7, TERTIARY, bold=True)
txt(s4, '24:37', 0.75, 5.1,  5.2, 1.5, 52, PRIMARY, bold=True)
box(s4, 0.88, 7.15, 5.1, 0.05, SURFACE_HHI)
box(s4, 0.88, 7.15, 2.5, 0.05, PRIMARY)
txt(s4, '00:23', 0.88, 7.28, 1.5, 0.3, 8, ON_SURFACE)
txt(s4, '25:00', 5.4,  7.28, 0.7, 0.3, 8, ON_SURFACE, align=PP_ALIGN.RIGHT)

# Input Fields
txt(s4, 'INPUT FIELDS', 7.5, 4.15, 6, 0.35, 8, ON_SURFACE, bold=True)
box(s4, 7.5, 4.55, 6.0, 0.62, SURFACE_LOW, None)
txt(s4, 'Placeholder text...', 7.7, 4.62, 5.6, 0.45, 11, OUTLINE)
box(s4, 7.5, 5.42, 6.0, 0.62, SURFACE_LOW, OUTLINE)
txt(s4, 'Active input field', 7.7, 5.49, 5.6, 0.45, 11, PRIMARY)

# Ghost border callout
txt(s4, 'Ghost Border  —  outline-variant at 15-20% opacity\nNever use 100% opaque outlines.', 7.5, 6.3, 7.5, 0.8, 9, ON_SURFACE)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Elevation
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)

badge(s5, 'FOUNDATIONS', 0.6, 0.4)
txt(s5, 'Elevation & Depth', 0.6, 0.78, 10, 0.8, 36, PRIMARY, bold=True)
txt(s5, 'No shadows — hierarchy defined entirely by tonal surface shifts', 0.6, 1.52, 12, 0.4, 12, ON_SURFACE)

layers = [
    (BG,          'Layer 0',  'surface',                    '#131313'),
    (SURFACE_LOW, 'Layer 1',  'surface-container-low',      '#1b1b1b'),
    (SURFACE,     'Layer 2',  'surface-container',          '#222222'),
    (SURFACE_HI,  'Layer 3',  'surface-container-high',     '#2a2a2a'),
    (SURFACE_HHI, 'Layer 4',  'surface-container-highest',  '#353535'),
]

for i, (color, layer, token, hex_val) in enumerate(layers):
    lx = 0.5 + i * 3.05
    h  = 2.2 + i * 0.55
    box(s5, lx, 8.6 - h, 2.75, h, color, OUTLINE)
    txt(s5, layer,   lx+0.12, 8.6 - h + 0.12, 2.4, 0.35, 8,  TERTIARY, bold=True)
    txt(s5, token,   lx+0.12, 8.6 - h + 0.48, 2.5, 0.35, 7,  ON_SURFACE)
    txt(s5, hex_val, lx+0.12, 8.6 - h + 0.80, 2.5, 0.35, 9,  PRIMARY,   bold=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Do's and Don'ts
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)

badge(s6, 'GUIDELINES', 0.6, 0.4)
txt(s6, "Do's & Don'ts", 0.6, 0.78, 12, 0.8, 36, PRIMARY, bold=True)

box(s6, 0.5, 1.82, 7.3, 0.06, TERTIARY)
txt(s6, 'DO', 0.6, 2.0, 3, 0.45, 13, TERTIARY, bold=True)
dos = [
    'Embrace negative space — increase padding before adding borders.',
    'Use tertiary (#6BFDAF) sparingly as a high-tech accent color.',
    'Apply backdrop-blur to all floating / overlapping elements.',
    'Define hierarchy through tonal layers, never border lines.',
    'Minimum border-radius of 0.25rem (sm) on all elements.',
]
for j, d in enumerate(dos):
    txt(s6, f'   {d}', 0.62, 2.6 + j * 0.78, 7.0, 0.65, 11, ON_SURFACE)

box(s6, 8.2, 1.82, 7.3, 0.06, BREAK_ACC)
txt(s6, "DON'T", 8.3, 2.0, 4, 0.45, 13, BREAK_ACC, bold=True)
donts = [
    'Use pure white for long-form body copy — causes halation.',
    'Use standard blue for links. Use primary (white) or tertiary.',
    'Use sharp 90-degree corners anywhere in the interface.',
    'Use 1px solid borders for structural layout separation.',
    'Use full-opacity outlines — ghost borders at 15-20% only.',
]
for j, d in enumerate(donts):
    txt(s6, f'   {d}', 8.32, 2.6 + j * 0.78, 7.0, 0.65, 11, ON_SURFACE)

# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out = r'C:/Users/Kendrick/Desktop/Repos/personal-proj/jax-timer/assets/obsidian-hud-design-system.pptx'
prs.save(out)
print('Saved:', out)
